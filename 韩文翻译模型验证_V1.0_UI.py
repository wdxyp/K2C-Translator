import tkinter as tk
from tkinter import filedialog, messagebox
import torch
import torch.nn as nn
import pickle
import os
import copy
from konlpy.tag import Okt
from openpyxl import load_workbook
from openpyxl.styles import Alignment
from openpyxl.cell.cell import MergedCell
from docx import Document
from pptx import Presentation
import numpy as np
import pandas as pd
from datetime import datetime

# ========== 1. 模型架构定义 (必须与训练代码完全一致) ==========
class Encoder(nn.Module):
    def __init__(self, input_dim, emb_dim, hid_dim, n_layers, dropout, batch_first=True):
        super().__init__()
        self.batch_first = batch_first
        self.embedding = nn.Embedding(input_dim, emb_dim)
        self.rnn = nn.LSTM(emb_dim, hid_dim, n_layers, dropout=dropout, batch_first=batch_first)
        self.dropout = nn.Dropout(dropout)

    def forward(self, src, src_lens):
        embedded = self.dropout(self.embedding(src))
        if isinstance(src_lens, torch.Tensor):
            src_lens = src_lens.to('cpu')
        packed = torch.nn.utils.rnn.pack_padded_sequence(
            embedded,
            src_lens,
            batch_first=self.batch_first,
            enforce_sorted=False,
        )
        _, (hidden, cell) = self.rnn(packed)
        return hidden, cell

class Decoder(nn.Module):
    def __init__(self, output_dim, emb_dim, hid_dim, n_layers, dropout, batch_first=True):
        super().__init__()
        self.output_dim = output_dim
        self.batch_first = batch_first
        self.embedding = nn.Embedding(output_dim, emb_dim)
        self.rnn = nn.LSTM(emb_dim, hid_dim, n_layers, dropout=dropout, batch_first=batch_first)
        self.fc_out = nn.Linear(hid_dim, output_dim)
        self.dropout = nn.Dropout(dropout)

    def forward(self, input, hidden, cell):
        input = input.unsqueeze(1) if self.batch_first else input.unsqueeze(0)
        embedded = self.dropout(self.embedding(input))
        output, (hidden, cell) = self.rnn(embedded, (hidden, cell))
        prediction = self.fc_out(output.squeeze(1) if self.batch_first else output.squeeze(0))
        return prediction, hidden, cell

class Seq2Seq(nn.Module):
    def __init__(self, encoder, decoder, device):
        super().__init__()
        self.encoder = encoder
        self.decoder = decoder
        self.device = device

    def forward(self, src, src_lens, trg, teacher_forcing_ratio=0.5):
        batch_size = src.shape[0]
        trg_len = trg.shape[1]
        outputs = torch.zeros(batch_size, trg_len, self.decoder.output_dim).to(src.device)
        
        hidden, cell = self.encoder(src, src_lens)
        input = trg[:, 0]
        
        for t in range(1, trg_len):
            output, hidden, cell = self.decoder(input, hidden, cell)
            outputs[:, t, :] = output
            teacher_force = np.random.random() < teacher_forcing_ratio
            input = trg[:, t] if teacher_force else output.argmax(1)
            
        return outputs

# ========== 2. 翻译应用逻辑 ==========
class TranslatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("韩语翻译模型验证工具 V1.0 (本地模型版)")
        self.root.geometry("900x850")
        
        # 设备检测
        self.device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
        self.model = None
        self._loaded_state_dict = None
        self._loaded_hparams = None
        self._infer_batch_first = True
        self._last_display_batch_first = None
        self.k_vocab = None
        self.c_vocab = None
        self.inv_c_vocab = None
        self.okt = Okt()
        
        # 语料收集变量
        self.original_texts = []
        self.translated_texts = []
        self.correction_map = {}
        
        # UI 变量
        self.model_path = tk.StringVar()
        self.is_unified_vocab = tk.BooleanVar(value=True)
        self.k_vocab_path = tk.StringVar()
        self.c_vocab_path = tk.StringVar()
        self.input_file_path = tk.StringVar()
        self.output_folder_path = tk.StringVar()
        self.custom_filename = tk.StringVar()
        self.corpus_file_path = tk.StringVar(value='Corpus.xlsx')
        self.generate_corpus = tk.BooleanVar(value=True)
        self.append_translation = tk.BooleanVar(value=False)
        
        self.create_widgets()

    def create_widgets(self):
        # 1. 模型加载区
        load_frame = tk.LabelFrame(self.root, text="第一步：加载本地模型权重与词汇表", padx=10, pady=10, fg="blue")
        load_frame.pack(fill="x", padx=15, pady=5)
        
        # 模型文件
        f1 = tk.Frame(load_frame)
        f1.pack(fill="x", pady=2)
        tk.Label(f1, text="模型权重 (.pth):", width=15, anchor="w").pack(side="left")
        tk.Entry(f1, textvariable=self.model_path).pack(side="left", fill="x", expand=True, padx=5)
        tk.Button(f1, text="浏览", command=lambda: self.browse_file(self.model_path, [("PyTorch Model", "*.pth")])).pack(side="right")
        
        # 词汇表模式切换
        tk.Checkbutton(load_frame, text="使用统合词汇表 (韩文/中文在同一个文件内)", 
                       variable=self.is_unified_vocab, command=self.toggle_vocab_ui).pack(anchor="w", pady=2)

        # 韩语词汇表 (统合模式下作为唯一入口)
        self.vocab_frame_k = tk.Frame(load_frame)
        self.vocab_frame_k.pack(fill="x", pady=2)
        self.lbl_k_vocab = tk.Label(self.vocab_frame_k, text="统合词汇表 (.pkl):", width=15, anchor="w")
        self.lbl_k_vocab.pack(side="left")
        tk.Entry(self.vocab_frame_k, textvariable=self.k_vocab_path).pack(side="left", fill="x", expand=True, padx=5)
        tk.Button(self.vocab_frame_k, text="浏览", command=lambda: self.browse_file(self.k_vocab_path, [("Pickle", "*.pkl")])).pack(side="right")
        
        # 中文词汇表 (仅在非统合模式下显示)
        self.vocab_frame_c = tk.Frame(load_frame)
        if not self.is_unified_vocab.get():
            self.vocab_frame_c.pack(fill="x", pady=2)
            
        tk.Label(self.vocab_frame_c, text="中文词汇 (.pkl):", width=15, anchor="w").pack(side="left")
        tk.Entry(self.vocab_frame_c, textvariable=self.c_vocab_path).pack(side="left", fill="x", expand=True, padx=5)
        tk.Button(self.vocab_frame_c, text="浏览", command=lambda: self.browse_file(self.c_vocab_path, [("Pickle", "*.pkl")])).pack(side="right")
        
        tk.Button(load_frame, text="🚀 点击加载模型", command=self.load_model_and_vocab, bg="#4CAF50", fg="white", font=("Arial", 10, "bold")).pack(pady=10)

        # 2. 文本翻译区
        text_frame = tk.LabelFrame(self.root, text="第二步：即时翻译测试 (韩 -> 中)", padx=10, pady=10)
        text_frame.pack(fill="x", padx=15, pady=5)
        
        tk.Label(text_frame, text="请输入韩语:").pack(anchor="w")
        self.input_text = tk.Text(text_frame, height=4)
        self.input_text.pack(fill="x", pady=5)
        
        tk.Button(text_frame, text="开始翻译 ↓", command=self.translate_text_ui).pack()
        
        tk.Label(text_frame, text="中文翻译结果:").pack(anchor="w")
        self.output_text = tk.Text(text_frame, height=4, bg="#f9f9f9")
        self.output_text.pack(fill="x", pady=5)

        # 3. 文件处理区
        file_frame = tk.LabelFrame(self.root, text="第三步：批量文件翻译 (保留格式)", padx=10, pady=10, fg="darkgreen")
        file_frame.pack(fill="x", padx=15, pady=5)
        
        # 输入文件
        f4 = tk.Frame(file_frame)
        f4.pack(fill="x", pady=2)
        tk.Label(f4, text="选择文件:", width=15, anchor="w").pack(side="left")
        tk.Entry(f4, textvariable=self.input_file_path).pack(side="left", fill="x", expand=True, padx=5)
        tk.Button(f4, text="浏览", command=lambda: self.browse_file(self.input_file_path, [("Office文件", "*.docx;*.xlsx;*.xls;*.pptx")])).pack(side="right")
        
        # 输出目录
        f5 = tk.Frame(file_frame)
        f5.pack(fill="x", pady=2)
        tk.Label(f5, text="保存目录:", width=15, anchor="w").pack(side="left")
        tk.Entry(f5, textvariable=self.output_folder_path).pack(side="left", fill="x", expand=True, padx=5)
        tk.Button(f5, text="选择目录", command=self.browse_directory).pack(side="right")
        
        # 自定义文件名
        f6 = tk.Frame(file_frame)
        f6.pack(fill="x", pady=2)
        tk.Label(f6, text="自定义文件名:", width=15, anchor="w").pack(side="left")
        tk.Entry(f6, textvariable=self.custom_filename).pack(side="left", fill="x", expand=True, padx=5)
        tk.Label(f6, text="(可选)").pack(side="right")
        
        # 纠错语料库
        f7 = tk.Frame(file_frame)
        f7.pack(fill="x", pady=2)
        tk.Label(f7, text="纠错语料库:", width=15, anchor="w").pack(side="left")
        tk.Entry(f7, textvariable=self.corpus_file_path).pack(side="left", fill="x", expand=True, padx=5)
        tk.Button(f7, text="浏览", command=lambda: self.browse_file(self.corpus_file_path, [("Excel文件", "*.xlsx;*.xls")])).pack(side="right")

        # 选项
        opt_frame = tk.Frame(file_frame)
        opt_frame.pack(fill="x", pady=5)
        tk.Checkbutton(opt_frame, text="保留原文 (在原文下方添加翻译文本)", variable=self.append_translation).pack(side="left", padx=10)
        tk.Checkbutton(opt_frame, text="生成新语料库", variable=self.generate_corpus).pack(side="left", padx=10)
        
        tk.Button(file_frame, text="🔥 开始批量翻译并保存", command=self.process_file, bg="#2196F3", fg="white", font=("Arial", 11, "bold")).pack(pady=10)
        
        # 底部状态栏
        self.status_label = tk.Label(self.root, text="状态: 等待模型加载...", bd=1, relief="sunken", anchor="w", bg="#eeeeee")
        self.status_label.pack(side="bottom", fill="x")

    def toggle_vocab_ui(self):
        if self.is_unified_vocab.get():
            self.vocab_frame_c.pack_forget()
            self.lbl_k_vocab.config(text="统合词汇表 (.pkl):")
        else:
            self.vocab_frame_c.pack(fill="x", pady=2)
            self.lbl_k_vocab.config(text="韩语词汇 (.pkl):")

    def browse_file(self, var, types):
        filename = filedialog.askopenfilename(filetypes=types)
        if filename:
            var.set(filename)

    def browse_directory(self):
        directory = filedialog.askdirectory()
        if directory:
            self.output_folder_path.set(directory)

    def _extract_state_dict(self, checkpoint):
        if isinstance(checkpoint, dict):
            for key in ("state_dict", "model_state_dict", "model", "net"):
                if key in checkpoint and isinstance(checkpoint[key], dict):
                    return checkpoint[key]
        return checkpoint

    def _infer_model_hparams(self, state_dict):
        emb_dim = None
        hid_dim = None
        encoder_layers = set()
        decoder_layers = set()

        enc_emb = state_dict.get("encoder.embedding.weight")
        if enc_emb is not None and hasattr(enc_emb, "shape") and len(enc_emb.shape) == 2:
            emb_dim = int(enc_emb.shape[1])

        enc_hh = state_dict.get("encoder.rnn.weight_hh_l0")
        if enc_hh is not None and hasattr(enc_hh, "shape") and len(enc_hh.shape) == 2:
            hid_dim = int(enc_hh.shape[1])

        for k in state_dict.keys():
            if k.startswith("encoder.rnn.weight_ih_l"):
                try:
                    encoder_layers.add(int(k.split("encoder.rnn.weight_ih_l", 1)[1].split(".", 1)[0]))
                except Exception:
                    pass
            elif k.startswith("decoder.rnn.weight_ih_l"):
                try:
                    decoder_layers.add(int(k.split("decoder.rnn.weight_ih_l", 1)[1].split(".", 1)[0]))
                except Exception:
                    pass

        n_layers = 2
        if encoder_layers:
            n_layers = max(encoder_layers) + 1
        if decoder_layers:
            n_layers = max(n_layers, max(decoder_layers) + 1)

        if emb_dim is None or hid_dim is None:
            raise ValueError("无法从模型权重推断 EMB_DIM/HID_DIM，请确认权重文件是否为本项目训练出的 Seq2Seq")

        return emb_dim, hid_dim, n_layers

    def _build_model(self, emb_dim, hid_dim, n_layers, batch_first):
        if self.k_vocab is None or self.c_vocab is None:
            raise ValueError("词汇表未加载")
        enc = Encoder(len(self.k_vocab), emb_dim, hid_dim, n_layers, 0.5, batch_first=batch_first)
        dec = Decoder(len(self.c_vocab), emb_dim, hid_dim, n_layers, 0.5, batch_first=batch_first)
        self.model = Seq2Seq(enc, dec, self.device).to(self.device)
        self._infer_batch_first = batch_first

    def _update_mode_status(self):
        if not self._loaded_hparams:
            return
        if self._last_display_batch_first == self._infer_batch_first:
            return
        emb_dim, hid_dim, n_layers = self._loaded_hparams
        mode = "batch_first=True" if self._infer_batch_first else "batch_first=False"
        self.status_label.config(text=f"✅ 模型就绪 (EMB={emb_dim}, HID={hid_dim}, LAYERS={n_layers}, {mode}, 运行设备: {self.device})")
        self._last_display_batch_first = self._infer_batch_first

    def _decode_greedy(self, sentence, max_len, batch_first):
        if not self.model or self.k_vocab is None or self.c_vocab is None or self.inv_c_vocab is None:
            return "", 1.0

        sentence = str(sentence).strip()
        tokens = self.okt.morphs(sentence)
        tokens = ['<sos>'] + tokens + ['<eos>']

        k_vocab = self.k_vocab
        c_vocab = self.c_vocab
        inv_c_vocab = self.inv_c_vocab

        unk_idx = k_vocab.get('<unk>', 3)
        src_indices = [k_vocab.get(token, unk_idx) for token in tokens]
        seq_len = len(src_indices)

        if batch_first:
            src_tensor = torch.LongTensor(src_indices).unsqueeze(0).to(self.device)
        else:
            src_tensor = torch.LongTensor(src_indices).unsqueeze(1).to(self.device)

        src_len = [seq_len]

        with torch.no_grad():
            hidden, cell = self.model.encoder(src_tensor, src_len)

        sos_idx = c_vocab.get('<sos>', 1)
        eos_idx = c_vocab.get('<eos>', 2)
        trg_indices = [sos_idx]

        for _ in range(max_len):
            trg_tensor = torch.LongTensor([trg_indices[-1]]).to(self.device)
            with torch.no_grad():
                output, hidden, cell = self.model.decoder(trg_tensor, hidden, cell)
            pred_token = output.argmax(1).item()
            if pred_token == eos_idx:
                break
            trg_indices.append(pred_token)

        translated_tokens = [inv_c_vocab.get(idx, '<unk>') for idx in trg_indices[1:]]
        if not translated_tokens:
            return "", 1.0

        unk_count = sum(1 for t in translated_tokens if t == '<unk>')
        unk_ratio = unk_count / max(1, len(translated_tokens))
        return "".join(translated_tokens), unk_ratio

    def load_model_and_vocab(self):
        try:
            m_path = self.model_path.get()
            kv_path = self.k_vocab_path.get()
            cv_path = self.c_vocab_path.get()
            
            # 校验输入
            if self.is_unified_vocab.get():
                if not (m_path and kv_path):
                    messagebox.showwarning("提示", "请先选择模型和统合词汇表文件！")
                    return
            else:
                if not (m_path and kv_path and cv_path):
                    messagebox.showwarning("提示", "请先选择模型、韩语词汇表和中文词汇表！")
                    return

            self.status_label.config(text="正在加载词汇表...")
            self.root.update()
            
            if self.is_unified_vocab.get():
                with open(kv_path, 'rb') as f:
                    combined_vocab = pickle.load(f)
                    # 支持列表/元组格式: [k_vocab, c_vocab]
                    if isinstance(combined_vocab, (list, tuple)) and len(combined_vocab) >= 2:
                        self.k_vocab = combined_vocab[0]
                        self.c_vocab = combined_vocab[1]
                    # 支持字典格式: {'ko': korean_vocab, 'zh': chinese_vocab}
                    elif isinstance(combined_vocab, dict):
                        if 'ko' in combined_vocab and 'zh' in combined_vocab:
                            self.k_vocab = combined_vocab['ko']
                            self.c_vocab = combined_vocab['zh']
                        elif 'korean' in combined_vocab and 'chinese' in combined_vocab:
                            self.k_vocab = combined_vocab['korean']
                            self.c_vocab = combined_vocab['chinese']
                        else:
                            raise ValueError("统合词汇表(字典)格式错误，缺少 'ko'/'zh' 或 'korean'/'chinese' 键")
                    else:
                        raise ValueError("统合词汇表格式不支持，请确保是 [ko, zh] 列表或 {'ko':.., 'zh':..} 字典")
            else:
                with open(kv_path, 'rb') as f: self.k_vocab = pickle.load(f)
                with open(cv_path, 'rb') as f: self.c_vocab = pickle.load(f)
            
            self.inv_c_vocab = {v: k for k, v in self.c_vocab.items()}
            
            self.status_label.config(text="正在构建模型架构...")
            self.root.update()
            
            self.status_label.config(text="正在加载模型权重...")
            self.root.update()
            
            checkpoint = torch.load(m_path, map_location=self.device)
            state_dict = self._extract_state_dict(checkpoint)
            emb_dim, hid_dim, n_layers = self._infer_model_hparams(state_dict)

            self._loaded_state_dict = state_dict
            self._loaded_hparams = (emb_dim, hid_dim, n_layers)
            self._build_model(emb_dim, hid_dim, n_layers, batch_first=True)
            if self.model is None:
                raise ValueError("模型构建失败")
            self.model.load_state_dict(state_dict)
            self.model.eval()
            
            self._update_mode_status()
            messagebox.showinfo("成功", "模型及词汇表加载成功，现在可以开始翻译了！")
        except Exception as e:
            self.status_label.config(text="❌ 加载失败")
            messagebox.showerror("加载失败", f"错误详情: {str(e)}")

    def load_correction_corpus(self):
        self.correction_map = {}
        corpus_path = self.corpus_file_path.get()
        if os.path.exists(corpus_path):
            try:
                self.status_label.config(text=f"正在加载纠错语料库 ({os.path.basename(corpus_path)})...")
                self.root.update()
                df = pd.read_excel(corpus_path)
                if '翻译后' in df.columns and '修改后' in df.columns:
                    for _, row in df.iterrows():
                        trans = str(row['翻译后']).strip()
                        mod = str(row['修改后']).strip()
                        if trans and mod:
                            self.correction_map[trans] = mod
                print(f"成功加载纠错语料库，共 {len(self.correction_map)} 条规则。")
            except Exception as e:
                print(f"加载纠错语料库失败: {e}")

    def apply_corrections(self, text):
        if not text: return text
        stripped = text.strip()
        return self.correction_map.get(stripped, text)

    def translate_sentence(self, sentence, max_len=100):
        if not self.model or self.k_vocab is None or self.c_vocab is None or self.inv_c_vocab is None:
            return "模型或词汇表未加载"
        if not sentence or not str(sentence).strip(): return ""

        sentence = str(sentence).strip()
        result, unk_ratio = self._decode_greedy(sentence, max_len, batch_first=self._infer_batch_first)
        if (not result or unk_ratio > 0.6) and self._loaded_state_dict and self._loaded_hparams:
            emb_dim, hid_dim, n_layers = self._loaded_hparams
            self._build_model(emb_dim, hid_dim, n_layers, batch_first=False)
            if self.model is None:
                return result
            self.model.load_state_dict(self._loaded_state_dict)
            self.model.eval()
            alt_result, alt_unk_ratio = self._decode_greedy(sentence, max_len, batch_first=False)
            if alt_result and alt_unk_ratio < unk_ratio:
                result = alt_result
                self._update_mode_status()
            else:
                emb_dim, hid_dim, n_layers = self._loaded_hparams
                self._build_model(emb_dim, hid_dim, n_layers, batch_first=True)
                if self.model is None:
                    return result
                self.model.load_state_dict(self._loaded_state_dict)
                self.model.eval()
                self._update_mode_status()
        
        # 应用纠错
        result = self.apply_corrections(result)
        
        # 记录用于语料库
        self.original_texts.append(sentence)
        self.translated_texts.append(result)
        
        return result

    def translate_text_ui(self):
        text = self.input_text.get("1.0", "end-1c").strip()
        if text:
            self.load_correction_corpus() # 单次翻译也加载纠错
            result = self.translate_sentence(text)
            self.output_text.delete("1.0", "end")
            self.output_text.insert("1.0", result)

    def append_logic(self, original, translated):
        if self.append_translation.get():
            original = original.strip()
            translated = translated.strip()
            if original and translated:
                return f"{original}\n{translated}"
            return original or translated
        return translated

    def process_file(self):
        input_path = self.input_file_path.get()
        output_dir = self.output_folder_path.get()
        
        if not input_path or not os.path.exists(input_path):
            messagebox.showerror("错误", "请选择有效的待翻译文件！")
            return
        if not output_dir:
            messagebox.showerror("错误", "请选择保存目录！")
            return
        if not self.model:
            messagebox.showerror("错误", "请先加载模型！")
            return

        self.original_texts = []
        self.translated_texts = []
        self.load_correction_corpus()

        file_ext = os.path.splitext(input_path)[1].lower()
        custom_name = self.custom_filename.get().strip()
        
        if custom_name:
            save_name = custom_name + file_ext
        else:
            save_name = "translated_" + os.path.basename(input_path)
            
        output_path = os.path.join(output_dir, save_name)

        try:
            self.status_label.config(text=f"🚀 正在处理: {os.path.basename(input_path)}...")
            self.root.update()
            
            if file_ext == '.docx':
                self.translate_docx(input_path, output_path)
            elif file_ext in ['.xlsx', '.xls']:
                self.translate_xlsx(input_path, output_path)
            elif file_ext == '.pptx':
                self.translate_pptx(input_path, output_path)
            else:
                messagebox.showerror("错误", "目前仅支持 docx, xlsx, xls, pptx 格式")
                return
            
            if self.generate_corpus.get():
                self.save_to_corpus()

            self.status_label.config(text=f"✅ 翻译完成！保存至: {os.path.basename(output_path)}")
            messagebox.showinfo("大功告成", f"文件已保存至:\n{output_path}")
        except Exception as e:
            self.status_label.config(text="❌ 翻译过程中断")
            messagebox.showerror("翻译失败", f"原因: {str(e)}")

    def save_to_corpus(self):
        if not self.original_texts: return
        try:
            corpus_folder = 'Corpus'
            if not os.path.exists(corpus_folder):
                os.makedirs(corpus_folder)
            timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
            file_name = f"Corpus_KO2ZH_{timestamp}.xlsx"
            path = os.path.join(corpus_folder, file_name)
            
            df = pd.DataFrame({
                '原文(韩语)': self.original_texts,
                '翻译后(中文)': self.translated_texts,
                '修改后': [''] * len(self.original_texts)
            })
            df.to_excel(path, index=False)
            print(f"语料库已生成: {path}")
        except Exception as e:
            print(f"生成语料库失败: {e}")

    # --- Word 处理 ---
    def translate_docx(self, input_path, output_path):
        doc = Document(input_path)
        
        # 1. 段落处理
        total_p = len(doc.paragraphs)
        for i, para in enumerate(doc.paragraphs):
            if i % 10 == 0:
                self.status_label.config(text=f"Word 进度: 段落 {i}/{total_p}...")
                self.root.update()
            
            for run in para.runs:
                # 先提取图片元素，防止 run.text = ... 清空它们
                drawings = run.element.xpath('.//w:drawing')
                
                # 处理文字翻译
                text = run.text
                if text.strip():
                    trans = self.translate_sentence(text)
                    run.text = self.append_logic(text, trans)
                
                # 重新添加图片元素
                for drawing in drawings:
                    run.element.append(drawing)

        # 2. 表格处理
        for i, table in enumerate(doc.tables):
            self.status_label.config(text=f"Word 进度: 表格 {i+1}/{len(doc.tables)}...")
            self.root.update()
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            drawings = run.element.xpath('.//w:drawing')
                            text = run.text
                            if text.strip():
                                trans = self.translate_sentence(text)
                                run.text = self.append_logic(text, trans)
                            
                            for drawing in drawings:
                                run.element.append(drawing)
        
        # 3. 形状文字处理 (参考 V2.9 XML 访问方式)
        for shape in doc.inline_shapes:
            try:
                if hasattr(shape, '_inline'):
                    inline = shape._inline
                    if hasattr(inline, 'graphic'):
                        graphic = inline.graphic
                        if hasattr(graphic, 'graphicData'):
                            graphic_data = graphic.graphicData
                            if hasattr(graphic_data, 'txBody'):
                                text_frame = graphic_data.txBody
                                if text_frame:
                                    # 遍历 XML 中的段落和运行
                                    for paragraph in text_frame.p:
                                        for run in paragraph.r:
                                            if hasattr(run, 't'):
                                                text = run.t
                                                if text and text.strip():
                                                    trans = self.translate_sentence(text)
                                                    run.t = self.append_logic(text, trans)
            except Exception:
                continue

        doc.save(output_path)

    # --- Excel 处理 ---
    def translate_xlsx(self, input_path, output_path):
        if input_path.lower().endswith('.xls'):
            # 使用 pandas 转换 xls 为 xlsx
            df_list = pd.read_excel(input_path, sheet_name=None)
            temp_path = input_path + "x"
            with pd.ExcelWriter(temp_path) as writer:
                for sheet_name, df in df_list.items():
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
            input_path = temp_path
            output_path = os.path.splitext(output_path)[0] + ".xlsx"

        wb = load_workbook(input_path)
        for sheet in wb.worksheets:
            max_r = sheet.max_row
            for row_idx, row in enumerate(sheet.iter_rows(), 1):
                if row_idx % 20 == 0:
                    self.status_label.config(text=f"Excel 进度: {sheet.title} 第 {row_idx}/{max_r} 行...")
                    self.root.update()
                for cell in row:
                    if isinstance(cell, MergedCell):
                        continue
                    if cell.value and isinstance(cell.value, str) and cell.value.strip():
                        # 备份格式
                        orig_font = copy.copy(cell.font)
                        orig_border = copy.copy(cell.border)
                        orig_fill = copy.copy(cell.fill)
                        
                        # 处理单元格内的换行
                        lines = str(cell.value).split('\n')
                        trans_lines = [self.translate_sentence(line) for line in lines]
                        trans_full = "\n".join(trans_lines)
                        
                        cell.value = self.append_logic(str(cell.value), trans_full)
                        
                        # 还原并优化格式
                        # 注意：在 openpyxl 中，直接赋值 StyleProxy 会触发类型检查错误，但运行时是可行的
                        # 使用 copy 确保获取的是具体样式对象
                        cell.font = orig_font # type: ignore
                        cell.border = orig_border # type: ignore
                        cell.alignment = Alignment(wrap_text=True, vertical='center')
                        cell.fill = orig_fill # type: ignore
                        
        wb.save(output_path)

    # --- PPT 处理 ---
    def translate_pptx(self, input_path, output_path):
        prs = Presentation(input_path)
        total_s = len(prs.slides)
        for i, slide in enumerate(prs.slides, 1):
            self.status_label.config(text=f"PPT 进度: 第 {i}/{total_s} 页...")
            self.root.update()
            for shape in slide.shapes:
                self.process_ppt_shape(shape)
        prs.save(output_path)

    def process_ppt_shape(self, shape):
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        # 备份格式
                        f_name = run.font.name
                        f_size = run.font.size
                        f_bold = run.font.bold
                        f_italic = run.font.italic
                        f_under = run.font.underline
                        f_color = run.font.color.rgb if run.font.color and hasattr(run.font.color, 'rgb') else None
                        
                        trans = self.translate_sentence(run.text)
                        run.text = self.append_logic(run.text, trans)
                        
                        # 还原格式
                        run.font.name = f_name
                        run.font.size = f_size
                        run.font.bold = f_bold
                        run.font.italic = f_italic
                        run.font.underline = f_under
                        if f_color: run.font.color.rgb = f_color
                        
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                trans = self.translate_sentence(run.text)
                                run.text = self.append_logic(run.text, trans)
                                
        elif shape.shape_type == 6: # 组合形状
            for s in shape.shapes:
                self.process_ppt_shape(s)

if __name__ == "__main__":
    root = tk.Tk()
    app = TranslatorApp(root)
    root.mainloop()
