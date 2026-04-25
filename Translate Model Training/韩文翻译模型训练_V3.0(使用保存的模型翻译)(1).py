# 翻译文件实操
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
import pandas as pd
import os
import random
import string
from openpyxl import load_workbook
from openpyxl.styles import Font, Border, Alignment, PatternFill
import torch
import torch.nn as nn
import pickle

# 假设这是之前定义的模型类； （定义编码器）
class Encoder(nn.Module):
    def __init__(self, input_dim, emb_dim, hid_dim, n_layers, dropout):
        super().__init__()
        self.hid_dim = hid_dim
        self.n_layers = n_layers
        self.embedding = nn.Embedding(input_dim, emb_dim)
        self.rnn = nn.LSTM(emb_dim, hid_dim, n_layers, dropout=dropout)
        self.dropout = nn.Dropout(dropout)

    def forward(self, src):
        embedded = self.dropout(self.embedding(src))
        outputs, (hidden, cell) = self.rnn(embedded)
        return hidden, cell

class Decoder(nn.Module):
    def __init__(self, output_dim, emb_dim, hid_dim, n_layers, dropout):
        super().__init__()
        self.output_dim = output_dim
        self.hid_dim = hid_dim
        self.n_layers = n_layers
        self.embedding = nn.Embedding(output_dim, emb_dim)
        self.rnn = nn.LSTM(emb_dim, hid_dim, n_layers, dropout=dropout)
        self.fc_out = nn.Linear(hid_dim, output_dim)
        self.dropout = nn.Dropout(dropout)

    def forward(self, input, hidden, cell):
        input = input.unsqueeze(0)
        embedded = self.dropout(self.embedding(input))
        output, (hidden, cell) = self.rnn(embedded, (hidden, cell))
        prediction = self.fc_out(output.squeeze(0))
        return prediction, hidden, cell

class Seq2Seq(nn.Module):
    def __init__(self, encoder, decoder, device):
        super().__init__()
        self.encoder = encoder
        self.decoder = decoder
        self.device = device

    def forward(self, src, trg, teacher_forcing_ratio = 0.5):
        batch_size = trg.shape[1]
        trg_len = trg.shape[0]
        trg_vocab_size = self.decoder.output_dim
        outputs = torch.zeros(trg_len, batch_size, trg_vocab_size).to(self.device)
        hidden, cell = self.encoder(src)
        input = trg[0,:]
        for t in range(1, trg_len):
            output, hidden, cell = self.decoder(input, hidden, cell)
            outputs[t] = output
            teacher_force = random.random() < teacher_forcing_ratio
            top1 = output.argmax(1)
            input = trg[t] if teacher_force else top1
        return outputs
# 加载模型参数
ENC_EMB_DIM = 256
DEC_EMB_DIM = 256
HID_DIM = 512
N_LAYERS = 2
ENC_DROPOUT = 0.5
DEC_DROPOUT = 0.5

# 加载训练好的模型
device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
# 这里需要根据实际情况修改输入维度、输出维度等参数
# 根据翻译方向动态初始化模型
def load_model(direction):
    model_file = model_file_entry.get()
    if not model_file:
        status_label.config(text="请选择模型文件。")
        return None
    # 检查词汇表是否加载
    if 'korean_vocab' not in globals() or 'chinese_vocab' not in globals():
        status_label.config(text="请先加载词汇表。")
        return None
    if not korean_vocab or not chinese_vocab:
        status_label.config(text="请先加载有效的词汇表。")
        return None
    # 设置维度参数
    if direction == 'zh2ko':
        input_dim = len(chinese_vocab)
        output_dim = len(korean_vocab)
    else:  # ko2zh
        input_dim = len(korean_vocab)
        output_dim = len(chinese_vocab)
    print(f"输入维度: {input_dim}, 输出维度: {output_dim}")
    encoder = Encoder(input_dim, ENC_EMB_DIM, HID_DIM, N_LAYERS, ENC_DROPOUT)
    decoder = Decoder(output_dim, DEC_EMB_DIM, HID_DIM, N_LAYERS, DEC_DROPOUT)
    model = Seq2Seq(encoder, decoder, device).to(device)
    try:
        model.load_state_dict(torch.load(model_file, map_location=device))
    except RuntimeError as e:
        status_label.config(text=f"加载模型参数时出错: {str(e)}")
        return None
    return model.eval()
    
def tokenize(text, direction):
    if direction == 'zh2ko':
        # 这里可以使用中文分词库，例如 jieba
        import jieba
        return list(jieba.cut(text))
    elif direction == 'ko2zh':
        # 对于韩文分词，可以使用 KoNLPy 库
        from konlpy.tag import Okt
        okt = Okt()
        return okt.morphs(text)
    return text.split()  # 简单示例，其他情况使用空格分词

# 修改后的翻译函数，支持三种语言互译
def get_translation(text, direction):
    model = load_model(direction)
    if model is None:
        return "无法加载模型，请选择有效的模型文件。"
    translated_text = ""
    tokens = tokenize(text, direction)
    print(f"分词结果: {tokens}")
    if direction == 'zh2ko':
        indices = [chinese_vocab.get(token, 0) for token in tokens]
        print(f"中文索引: {indices}")
        src = torch.tensor(indices).unsqueeze(1).to(device)
        trg = torch.zeros(10, 1, dtype=torch.long).to(device)
        with torch.no_grad():
            output = model(src, trg, 0)
        output = output.argmax(2)
        output_tokens = [korean_vocab.get(index.item(), '<unk>') for index in output.squeeze()]
        print(f"韩语输出索引: {output.tolist()}")
        print(f"韩语输出词: {output_tokens}")
        translated_text = ' '.join(output_tokens)
    elif direction == 'ko2zh':
        indices = [korean_vocab.get(token, 0) for token in tokens]
        print(f"韩语索引: {indices}")
        src = torch.tensor(indices).unsqueeze(1).to(device)
        trg = torch.zeros(10, 1, dtype=torch.long).to(device)
        with torch.no_grad():
            output = model(src, trg, 0)
        output = output.argmax(2)
        output_tokens = [chinese_vocab.get(index.item(), '<unk>') for index in output.squeeze()]
        print(f"中文输出索引: {output.tolist()}")
        print(f"中文输出词: {output_tokens}")
        translated_text = ' '.join(output_tokens)
    return translated_text

# 翻译TXT文件的函数
def translate_txt(input_file, output_file, direction):
    with open(input_file, 'r', encoding='utf-8') as f:
        text = f.read()
    translated_text = get_translation(text, direction)
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(translated_text)

# 翻译Word文件的函数
from docx import Document
def translate_word(input_file, output_file, direction):
    doc = Document(input_file)
    for para in doc.paragraphs:
        original_text = para.text
        translated_text = get_translation(original_text, direction)
        para.text = translated_text
    doc.save(output_file)

# 假设这里有 translate_ppt 和 translate_excel 函数的定义
def translate_ppt(input_file, output_file):
    prs = Presentation(input_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            # 双重安全检查：属性存在性检查 + 类型忽略注释
            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame  # type: ignore
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        direction = translation_direction.get()
                        translated_text = get_translation(original_text, direction)
                        run.text = translated_text
    prs.save(output_file)

def translate_excel(input_file, output_file):
    df = pd.read_excel(input_file)
    for col in df.columns:
        df[col] = df[col].apply(lambda x: get_translation(str(x), 'zh2ko'))  # 示例方向
    df.to_excel(output_file, index=False)

# 修改start_translation函数以支持新的文件类型
def start_translation():
    if 'korean_vocab' not in globals() or 'chinese_vocab' not in globals():
        status_label.config(text="请先加载词汇表。")
        return
    if not korean_vocab or not chinese_vocab:
        status_label.config(text="请先加载有效的词汇表。")
        return
    input_file = input_file_entry.get()
    output_folder = output_folder_entry.get()
    custom_filename = custom_filename_entry.get()
    direction = translation_direction.get()
    if input_file and output_folder:
        file_ext = os.path.splitext(input_file)[1]
        if custom_filename:
            output_file = os.path.join(output_folder, custom_filename + file_ext)
        else:
            file_name = os.path.basename(input_file)
            output_file = os.path.join(output_folder, file_name)
        if input_file.lower().endswith(('.ppt', '.pptx')):
            translate_ppt(input_file, output_file)
        elif input_file.lower().endswith(('.xls', '.xlsx')):
            translate_excel(input_file, output_file)
        elif input_file.lower().endswith('.txt'):
            translate_txt(input_file, output_file, direction)
        elif input_file.lower().endswith('.docx'):
            translate_word(input_file, output_file, direction)
        status_label.config(text="翻译完成！文件保存路径: " + output_file)
    else:
        status_label.config(text="请选择输入文件和输出文件夹路径。")

# 创建主窗口
root = tk.Tk()
root.title("文件翻译工具")

# 创建左右两个 Frame
left_frame = tk.Frame(root)
left_frame.pack(side=tk.LEFT, padx=10, pady=10)

right_frame = tk.Frame(root)
right_frame.pack(side=tk.RIGHT, padx=10, pady=10)

# 输入文件选择
input_file_label = tk.Label(left_frame, text="选择输入文件:")
input_file_label.pack(pady=5)
input_file_entry = tk.Entry(left_frame, width=50)
input_file_entry.pack(pady=2)
input_file_button = tk.Button(left_frame, text="浏览", command=lambda: input_file_entry.insert(0, filedialog.askopenfilename()))
input_file_button.pack(pady=5)

# 输出文件夹选择
output_folder_label = tk.Label(left_frame, text="选择输出文件夹:")
output_folder_label.pack(pady=5)
output_folder_entry = tk.Entry(left_frame, width=50)
output_folder_entry.pack(pady=2)
output_folder_button = tk.Button(left_frame, text="浏览", command=lambda: output_folder_entry.insert(0, filedialog.askdirectory()))
output_folder_button.pack(pady=5)

# 自定义文件名输入
custom_filename_label = tk.Label(left_frame, text="自定义输出文件名 (可选):")
custom_filename_label.pack(pady=5)
custom_filename_entry = tk.Entry(left_frame, width=50)
custom_filename_entry.pack(pady=2)

# 翻译方向选择
translation_direction = tk.StringVar()
translation_direction.set('zh2ko')

zh2ko_radio = tk.Radiobutton(left_frame, text="中文 -> 韩文", variable=translation_direction, value='zh2ko')
zh2ko_radio.pack(pady=2)

ko2zh_radio = tk.Radiobutton(left_frame, text="韩文 -> 中文", variable=translation_direction, value='ko2zh')
ko2zh_radio.pack(pady=2)

# 模型文件选择
model_file_label = tk.Label(left_frame, text="选择模型文件:")
model_file_label.pack(pady=5)
model_file_entry = tk.Entry(left_frame, width=50)
model_file_entry.pack(pady=2)
model_file_button = tk.Button(left_frame, text="浏览", command=lambda: model_file_entry.insert(0, filedialog.askopenfilename()))
model_file_button.pack(pady=5)

# 韩语词汇表文件选择
korean_vocab_file_label = tk.Label(right_frame, text="选择韩语词汇表文件:")
korean_vocab_file_label.pack(pady=5)
korean_vocab_file_entry = tk.Entry(right_frame, width=50)
korean_vocab_file_entry.pack(pady=2)
korean_vocab_file_button = tk.Button(right_frame, text="浏览", command=lambda: korean_vocab_file_entry.insert(0, filedialog.askopenfilename()))
korean_vocab_file_button.pack(pady=5)

# 中文词汇表文件选择
chinese_vocab_file_label = tk.Label(right_frame, text="选择中文词汇表文件:")
chinese_vocab_file_label.pack(pady=5)
chinese_vocab_file_entry = tk.Entry(right_frame, width=50)
chinese_vocab_file_entry.pack(pady=2)
chinese_vocab_file_button = tk.Button(right_frame, text="浏览", command=lambda: chinese_vocab_file_entry.insert(0, filedialog.askopenfilename()))
chinese_vocab_file_button.pack(pady=5)

# 修改加载词汇表部分
def load_vocab():
    global korean_vocab, chinese_vocab
    korean_vocab_file = korean_vocab_file_entry.get()
    chinese_vocab_file = chinese_vocab_file_entry.get()
    if korean_vocab_file and chinese_vocab_file:
        try:
            with open(korean_vocab_file, 'rb') as f:
                korean_vocab = pickle.load(f)
            with open(chinese_vocab_file, 'rb') as f:
                chinese_vocab = pickle.load(f)
            print(f"韩语词汇表大小: {len(korean_vocab)}")
            print(f"中文词汇表内容示例: {list(chinese_vocab.items())[:10]}")  # 打印前10个词汇查看
            print(f"中文词汇表大小: {len(chinese_vocab)}")
            status_label.config(text="词汇表加载成功！")

            # 显示韩语词汇表预览
            korean_preview_text = "韩语词汇表预览:\n"
            for word, idx in list(korean_vocab.items())[:10]:
                korean_preview_text += f"{word}: {idx}\n"
            korean_preview_label.config(text=korean_preview_text)

            # 显示中文词汇表预览
            chinese_preview_text = "中文词汇表预览:\n"
            for word, idx in list(chinese_vocab.items())[:10]:
                chinese_preview_text += f"{word}: {idx}\n"
            chinese_preview_label.config(text=chinese_preview_text)

        except Exception as e:
            status_label.config(text=f"加载词汇表时出错: {str(e)}")
    else:
        status_label.config(text="请选择韩语和中文词汇表文件。")

# 加载词汇表按钮
load_vocab_button = tk.Button(right_frame, text="加载词汇表", command=load_vocab)
load_vocab_button.pack(pady=5)

# 韩语词汇表预览标签
korean_preview_label = tk.Label(right_frame, text="", justify=tk.LEFT)
korean_preview_label.pack(pady=5)

# 中文词汇表预览标签
chinese_preview_label = tk.Label(right_frame, text="", justify=tk.LEFT)
chinese_preview_label.pack(pady=5)

# 开始翻译按钮
start_button = tk.Button(left_frame, text="开始翻译", command=start_translation)
start_button.pack(pady=20)

# 状态显示标签
status_label = tk.Label(left_frame, text="")
status_label.pack(pady=5)

# 运行主循环
root.mainloop()