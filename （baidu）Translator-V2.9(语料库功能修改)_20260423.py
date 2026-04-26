import tkinter as tk
from tkinter import filedialog
from PIL.ImageFont import truetype
from openpyxl.styles.alignment import vertical_aligments
from pptx import Presentation
import pandas as pd
import os
import hashlib
import random
import string
import requests
from openpyxl import load_workbook
from openpyxl.styles import Font, Border, Alignment, PatternFill
from docx import Document
from datetime import datetime
import re
import time

# 请替换为你自己的 APP ID 和密钥
APP_ID = '20250221002280429'
SECRET_KEY = '6kkWTs176PJjdGdPdy_Q'
API_URL = 'http://api.fanyi.baidu.com/api/trans/vip/translate'

# 创建主窗口
root = tk.Tk()
root.title("文件翻译工具")

# 全局变量用于存储是否生成语料库的选项
generate_corpus = tk.BooleanVar()
generate_corpus.set(False)  # 默认启用

# 全局变量用于存储是否在原文下方添加翻译文本的选项
append_translation = tk.BooleanVar()
append_translation.set(False)  # 默认不启用

# 全局变量用于存储纠错映射
correction_map = {}

# 全局变量用于存储语料库文件路径
corpus_file_path = tk.StringVar()
corpus_file_path.set('Corpus.xlsx')


def load_correction_corpus():
    """
    加载语料库文件进行纠错
    """
    global correction_map
    correction_map = {}
    corpus_path = corpus_file_path.get()
    if os.path.exists(corpus_path):
        try:
            msg = f"正在加载纠错语料库 ({corpus_path})..."
            status_label.config(text=msg)
            print(msg)
            root.update()
            df = pd.read_excel(corpus_path)
            if '翻译后' in df.columns and '修改后' in df.columns:
                for _, row in df.iterrows():
                    trans = str(row['翻译后']).strip()
                    mod = str(row['修改后']).strip()
                    if trans and mod:
                        correction_map[trans] = mod
            msg = f"成功加载纠错语料库，共 {len(correction_map)} 条规则。"
            print(msg)
        except Exception as e:
            msg = f"加载纠错语料库失败: {e}"
            print(msg)
            status_label.config(text=msg)
            root.update()
    else:
        print(f"未找到语料库文件 ({corpus_path})，将不进行内容替换。")


def apply_corrections(translated_text):
    """
    应用纠错映射，将“翻译后”的内容替换为“修改后”的内容
    """
    if not translated_text:
        return translated_text
    
    # 去除两端空格进行匹配
    stripped_text = translated_text.strip()
    if stripped_text in correction_map:
        return correction_map[stripped_text]
    return translated_text


def append_translation_to_original(text, translated_text, cell=None):
    """
    将翻译后的文字添加到原文下方，并处理单元格格式
    :param text: 原文
    :param translated_text: 翻译后的文本
    :param cell: 单元格对象，用于设置格式（仅用于.xlsx）
    :return: 添加翻译后的文本
    """
    text = text.strip()
    translated_text = translated_text.strip()
    if text and translated_text:
        result = f"{text}\n{translated_text}"
    else:
        result = text or translated_text
    
    if cell:
        # 设置单元格自动换行
        cell.alignment = Alignment(wrap_text=True,vertical='center')
        # 自动调整行高
        ws = cell.parent
        row_num = cell.row
        # 获取当前行的原始高度
        original_height = ws.row_dimensions[row_num].height
        # 计算新的行高为原文行高的2倍
        new_height = original_height * 2 if original_height else 0
        # 设置新的行高
        ws.row_dimensions[row_num].height = new_height
    return result

def get_translation(text):
    # 设置重试次数和延迟
    max_retries = 3
    retry_delay = 2  # 秒
    
    if translation_direction.get() == 'en2zh':
        # 按句子分割英文文本
        sentences = re.split(r'(?<=[.!?])\s+', text)
        translated_sentences = []
        for sentence in sentences:
            if not sentence:
                continue
            # 分割文本为数字、字母和文字片段，支持越南语字符
            pattern = r'([a-zA-Z0-9\u0100-\u017F]+|[^\u4e00-\u9fa5a-zA-Z0-9\u0100-\u017F]+)'
            segments = re.split(pattern, sentence)
            segment_translated = []
            from_lang = 'en'
            to_lang = 'zh'

            for segment in segments:
                segment = str(segment)
                if segment.isdigit():
                    segment_translated.append(segment)
                elif segment.isalpha() and segment.isascii() and from_lang != 'en':
                    segment_translated.append(segment)
                else:
                    if segment.strip():
                        # 添加重试逻辑
                        success = False
                        retries = 0
                        while not success and retries < max_retries:
                            try:
                                salt = ''.join(random.choices(string.ascii_letters + string.digits, k=16))
                                sign_str = APP_ID + segment + salt + SECRET_KEY
                                sign = hashlib.md5(sign_str.encode()).hexdigest()
                                params = {
                                    'q': segment,
                                    'from': from_lang,
                                    'to': to_lang,
                                    'appid': APP_ID,
                                    'salt': salt,
                                    'sign': sign
                                }
                                # 添加超时设置
                                response = requests.get(API_URL, params=params, timeout=10)
                                response.encoding = 'utf-8'
                                result = response.json()
                                if 'trans_result' in result:
                                    segment_translated.append(result['trans_result'][0]['dst'])
                                else:
                                    segment_translated.append(segment)
                                success = True
                            except requests.exceptions.RequestException as e:
                                retries += 1
                                print(f"翻译出错 (重试 {retries}/{max_retries}): {e}")
                                if retries < max_retries:
                                    time.sleep(retry_delay)
                                else:
                                    # 重试失败，使用原文
                                    segment_translated.append(segment)
                                    print(f"翻译失败，使用原文: {segment}")
                    else:
                        segment_translated.append(segment)

            translated_sentence = ''.join(segment_translated)
            translated_sentences.append(translated_sentence)

        return ' '.join(translated_sentences)
    elif translation_direction.get() == 'vi2zh':
        # 越南文整体翻译为中文，避免按片段拆分导致错误
        success = False
        retries = 0
        while not success and retries < max_retries:
            try:
                salt = ''.join(random.choices(string.ascii_letters + string.digits, k=16))
                sign_str = APP_ID + text + salt + SECRET_KEY
                sign = hashlib.md5(sign_str.encode()).hexdigest()
                params = {
                    'q': text,
                    'from': 'vie',
                    'to': 'zh',
                    'appid': APP_ID,
                    'salt': salt,
                    'sign': sign
                }
                response = requests.get(API_URL, params=params, timeout=10)
                response.encoding = 'utf-8'
                result = response.json()
                if 'trans_result' in result:
                    return result['trans_result'][0]['dst']
                else:
                    return text
            except requests.exceptions.RequestException as e:
                retries += 1
                print(f"翻译出错 (重试 {retries}/{max_retries}): {e}")
                if retries < max_retries:
                    time.sleep(retry_delay)
                else:
                    print(f"翻译失败，使用原文: {text}")
                    return text
    else:
        # 原有的翻译逻辑，支持越南语字符
        pattern = r'(\\w+|\\s+|[^\\w\\s]+)'
        segments = re.split(pattern, text)
        translated_segments = []
        from_lang = {
            'zh2ko': 'zh',
            'ko2zh': 'kor',
            'ko2en': 'kor',
            'zh2en': 'zh',
            'en2zh': 'en',
            'zh_tw2en': 'zh',
            'en2zh_tw': 'en',
            'zh2ja': 'zh',
            'ja2zh': 'ja',
            'en2ko': 'en',  # 新增英文到韩文的源语言映射
            'vi2zh': 'vie',  # 新增越南文到中文的源语言映射
            'zh2vi': 'zh',
            'ko2vi': 'kor'
        }.get(translation_direction.get(), 'auto')

        for segment in segments:
            segment = str(segment)
            if segment.isdigit():
                translated_segments.append(segment)
            elif segment.isalpha() and segment.isascii() and from_lang != 'en':
                translated_segments.append(segment)
            else:
                if segment.strip():
                    # 添加重试逻辑
                    success = False
                    retries = 0
                    while not success and retries < max_retries:
                        try:
                            salt = ''.join(random.choices(string.ascii_letters + string.digits, k=16))
                            sign_str = APP_ID + segment + salt + SECRET_KEY
                            sign = hashlib.md5(sign_str.encode()).hexdigest()
                            to_lang = {
                                'zh2ko': 'kor',
                                'ko2zh': 'zh',
                                'ko2en': 'en',
                                'zh2en': 'en',
                                'en2zh': 'zh',
                                'zh_tw2en': 'en',
                                'en2zh_tw': 'zh',
                                'zh2ja': 'ja',
                                'ja2zh': 'zh',
                                'en2ko': 'kor',  # 新增英文到韩文的目标语言映射
                                'vi2zh': 'zh',  # 新增越南文到中文的目标语言映射
                                'zh2vi': 'vie',
                                'ko2vi': 'vie'
                            }.get(translation_direction.get(), 'auto')
                            params = {
                                'q': segment,
                                'from': from_lang,
                                'to': to_lang,
                                'appid': APP_ID,
                                'salt': salt,
                                'sign': sign
                            }
                            # 添加超时设置
                            response = requests.get(API_URL, params=params, timeout=10)
                            response.encoding = 'utf-8'
                            result = response.json()
                            if 'trans_result' in result:
                                translated_segments.append(result['trans_result'][0]['dst'])
                            else:
                                translated_segments.append(segment)
                            success = True
                        except requests.exceptions.RequestException as e:
                            retries += 1
                            print(f"翻译出错 (重试 {retries}/{max_retries}): {e}")
                            if retries < max_retries:
                                time.sleep(retry_delay)
                            else:
                                # 重试失败，使用原文
                                translated_segments.append(segment)
                                print(f"翻译失败，使用原文: {segment}")
                else:
                    translated_segments.append(segment)

        translated_text = ''.join(translated_segments)
        return translated_text

# 全局变量用于存储翻译前后的文本
original_texts = []
translated_texts = []
def translate_shape_for_ppt(shape):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                # 保存原有格式
                original_font_name = run.font.name if run.font.name else None
                original_size = run.font.size
                original_bold = run.font.bold
                original_italic = run.font.italic
                original_underline = run.font.underline

                translated_text = get_translation(text)
                # 对翻译结果进行纠错替换
                translated_text = apply_corrections(translated_text)
                
                # 收集翻译前后的文本
                original_texts.append(text)
                translated_texts.append(translated_text)
                if append_translation.get():
                    run.text = append_translation_to_original(text, translated_text)
                else:
                    run.text = translated_text

                # 恢复原有格式
                if original_font_name is not None:
                    run.font.name = original_font_name
                run.font.size = original_size
                run.font.bold = original_bold
                run.font.italic = original_italic
                run.font.underline = original_underline
    elif shape.shape_type == 6:  # 组合形状
        for sub_shape in shape.shapes:
            translate_shape_for_ppt(sub_shape)
    elif shape.has_table:  # 处理表格形状
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        # 保存原有格式
                        original_font_name = run.font.name if run.font.name else None
                        original_size = run.font.size
                        original_bold = run.font.bold
                        original_italic = run.font.italic
                        original_underline = run.font.underline

                        translated_text = get_translation(text)
                        # 对翻译结果进行纠错替换
                        translated_text = apply_corrections(translated_text)
                        
                        # 收集翻译前后的文本
                        original_texts.append(text)
                        translated_texts.append(translated_text)
                        if append_translation.get():
                            run.text = append_translation_to_original(text, translated_text)
                        else:
                            run.text = translated_text

                        # 恢复原有格式
                        if original_font_name is not None:
                            run.font.name = original_font_name
                        run.font.size = original_size
                        run.font.bold = original_bold
                        run.font.italic = original_italic
                        run.font.underline = original_underline

def translate_ppt(input_file, output_file):
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        msg = f"正在翻译 PPT: 第 {i}/{total_slides} 页..."
        status_label.config(text=msg)
        print(msg)
        root.update()
        for shape in slide.shapes:
            # 直接调用 translate_shape_for_ppt 函数处理形状
            translate_shape_for_ppt(shape)
    
    msg = "正在保存 PPT 文件..."
    status_label.config(text=msg)
    print(msg)
    root.update()
    prs.save(output_file)
    # 翻译完成后调用保存语料库函数
    save_to_corpus(original_texts, translated_texts)
    # 清空列表
    original_texts.clear()
    translated_texts.clear()
    # status_label.config(text="PPT 文件翻译完成。")
    # root.update()
    print("保存修改后的 PPT 文件完成!")

def translate_excel(input_file, output_file):
    # 声明全局变量
    global status_label, original_texts, translated_texts
    
    # 检查文件扩展名，区分处理.xls和.xlsx文件
    file_ext = os.path.splitext(input_file)[1].lower()
    
    if file_ext == '.xlsx':
        translate_excel_xlsx(input_file, output_file)
        # 使用返回的新路径更新status_label
        if 'status_label' in globals():
            msg = "翻译完成！文件保存路径: " + output_file
            status_label.config(text=msg)
            #print(msg)
        # 翻译完成后调用保存语料库函数
        save_to_corpus(original_texts, translated_texts)
        # 清空列表
        original_texts.clear()
        translated_texts.clear()
    elif file_ext == '.xls':
        # 调用translate_excel_xls并获取返回的.xlsx文件路径
        output_file_xlsx = translate_excel_xls(input_file, output_file)
        # 使用返回的新路径更新status_label
        if 'status_label' in globals():
            msg = "翻译完成！文件保存路径: " + output_file_xlsx
            status_label.config(text=msg)
            #print(msg)
        # 翻译完成后调用保存语料库函数（在translate_excel_xls中已完成语料库的收集）
        save_to_corpus(original_texts, translated_texts)
        # 清空列表
        original_texts.clear()
        translated_texts.clear()
    else:
        raise ValueError(f"不支持的Excel格式: {file_ext}")

def clean_sheet_name(name):
    """清理和验证Excel工作表名称，确保符合Excel的命名规则"""
    if not name:
        return "Sheet"
    
    # 移除Excel不允许的所有特殊字符
    # Excel不允许的字符包括：\ / ? : * [ ] ( )
    invalid_chars = r'[\\/?:*\[\](){}<>|"\']'
    clean_name = re.sub(invalid_chars, '', name)
    
    # 确保名称长度不超过31个字符
    clean_name = clean_name[:31]
    
    # 确保名称不为空
    if not clean_name:
        clean_name = "Sheet"
    
    return clean_name

def translate_excel_xlsx(input_file, output_file):
    """处理.xlsx格式的Excel文件"""
    # 加载工作簿
    msg = "正在加载 Excel 文件..."
    status_label.config(text=msg)
    print(msg)
    root.update()
    wb = load_workbook(input_file)
    print(f"{input_file} 加载完成!")
    
    # 收集所有已存在的工作表名称（包括原始和已处理的）
    existing_sheet_names = set()
    
    total_sheets = len(wb.sheetnames)
    for i, sheet_name in enumerate(wb.sheetnames, 1):
        msg = f"正在翻译 Excel: 第 {i}/{total_sheets} 个工作表 ({sheet_name})..."
        status_label.config(text=msg)
        print(msg)
        root.update()
        # 不翻译工作表名称，只进行清理
        cleaned_sheet_name = clean_sheet_name(sheet_name)
        
        # 确保工作表名称唯一
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            temp_name = f"{cleaned_sheet_name[:27]}_{count}"
            unique_sheet_name = temp_name[:31]
            count += 1
        
        existing_sheet_names.add(unique_sheet_name)
        
        ws = wb[sheet_name]
        # 修改工作表名称
        ws.title = unique_sheet_name
        
        # 获取总行数用于进度（大致）
        total_rows = ws.max_row
        for row_idx, row in enumerate(ws.iter_rows(), 1):
            if row_idx % 10 == 0: # 每10行更新一次状态，避免过于频繁
                msg = f"正在翻译 Excel: {sheet_name} 第 {row_idx}/{total_rows} 行..."
                status_label.config(text=msg)
                print(msg)
                root.update()
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    # 按换行符分割文本
                    lines = cell.value.split('\n')
                    translated_lines = []
                    for line in lines:
                        translated_text = get_translation(line)
                        # 对翻译结果进行纠错替换
                        translated_text = apply_corrections(translated_text)
                        
                        # 收集翻译前后的文本
                        original_texts.append(line)
                        translated_texts.append(translated_text)
                        translated_lines.append(translated_text)
                    # 重新组合翻译后的行
                    translated_text = '\n'.join(translated_lines)
                    # 保存原格式
                    original_font = Font(**cell.font.__dict__)
                    original_border = Border(**cell.border.__dict__)
                    original_alignment = Alignment(**cell.alignment.__dict__)
                    original_fill = PatternFill(**cell.fill.__dict__)
                    if append_translation.get():
                        cell.value = append_translation_to_original(cell.value, translated_text, cell)
                    else:
                        cell.value = translated_text
                    # 恢复原格式
                    cell.font = original_font
                    cell.border = original_border
                    cell.alignment = original_alignment
                    cell.fill = original_fill
                    # 设置单元格自动换行
                    cell.alignment = Alignment(wrap_text=True,vertical='center')
    
    # 保存修改后的工作簿
    msg = "正在保存 Excel 文件..."
    status_label.config(text=msg)
    print(msg)
    root.update()
    wb.save(output_file)
    print("保存修改后的 Excel 文件完成!")

def translate_excel_xls(input_file, output_file):
    """处理.xls格式的Excel文件，使用pandas库，保存为.xlsx格式"""
    msg = "正在加载 Excel (.xls) 文件..."
    status_label.config(text=msg)
    print(msg)
    root.update()
    print(f"{input_file} 加载完成!")
    
    # 修改输出文件扩展名，保存为.xlsx格式
    base_output = os.path.splitext(output_file)[0]
    output_file_xlsx = base_output + '.xlsx'
    print(f"由于系统限制，.xls文件将翻译为.xlsx格式：{output_file_xlsx}")
    
    # 读取所有工作表
    excel_file = pd.ExcelFile(input_file)
    
    # 创建一个新的ExcelWriter对象，使用默认引擎
    writer = pd.ExcelWriter(output_file_xlsx, engine='openpyxl')
    
    # 收集所有已存在的工作表名称（包括原始和已处理的）
    existing_sheet_names = set()
    
    total_sheets = len(excel_file.sheet_names)
    for i, sheet_name in enumerate(excel_file.sheet_names, 1):
        msg = f"正在翻译 Excel: 第 {i}/{total_sheets} 个工作表 ({sheet_name})..."
        status_label.config(text=msg)
        print(msg)
        root.update()
        # 不翻译工作表名称，只进行清理
        cleaned_sheet_name = clean_sheet_name(sheet_name)
        
        # 确保工作表名称唯一
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            temp_name = f"{cleaned_sheet_name[:27]}_{count}"
            unique_sheet_name = temp_name[:31]
            count += 1
        
        existing_sheet_names.add(unique_sheet_name)
        
        # 读取工作表
        df = pd.read_excel(input_file, sheet_name=sheet_name)
        
        # 遍历所有单元格进行翻译
        total_rows = len(df.index)
        for idx_num, idx in enumerate(df.index, 1):
            if idx_num % 10 == 0:
                msg = f"正在翻译 Excel: {sheet_name} 第 {idx_num}/{total_rows} 行..."
                status_label.config(text=msg)
                print(msg)
                root.update()
            for col in df.columns:
                cell_value = df.at[idx, col]
                if pd.notna(cell_value) and isinstance(cell_value, str):
                    # 按换行符分割文本
                    lines = cell_value.split('\n')
                    translated_lines = []
                    for line in lines:
                        translated_text = get_translation(line)
                        # 对翻译结果进行纠错替换
                        translated_text = apply_corrections(translated_text)
                        
                        # 收集翻译前后的文本
                        original_texts.append(line)
                        translated_texts.append(translated_text)
                        translated_lines.append(translated_text)
                    # 重新组合翻译后的行
                    translated_text = '\n'.join(translated_lines)
                    
                    if append_translation.get():
                        df.at[idx, col] = append_translation_to_original(cell_value, translated_text)
                    else:
                        df.at[idx, col] = translated_text
        
        # 保存翻译后的工作表
        df.to_excel(writer, sheet_name=unique_sheet_name, index=False)
    
    # 保存工作簿
    msg = "正在保存 Excel 文件..."
    status_label.config(text=msg)
    print(msg)
    root.update()
    writer.close()
    
    # 返回新的.xlsx文件路径
    return output_file_xlsx

def translate_word(input_file, output_file):
    msg = "正在加载 Word 文件..."
    status_label.config(text=msg)
    print(msg)
    root.update()
    doc = Document(input_file)
    
    # 处理普通段落
    total_paragraphs = len(doc.paragraphs)
    for i, paragraph in enumerate(doc.paragraphs, 1):
        if i % 5 == 0:
            msg = f"正在翻译 Word: 段落 {i}/{total_paragraphs}..."
            status_label.config(text=msg)
            print(msg)
            root.update()
        for run in paragraph.runs:
            # 处理文字翻译
            text = run.text
            if text.strip():  # 只翻译非空文本
                # 先获取翻译结果
                translated_text = get_translation(text)
                # 对翻译结果进行纠错替换
                translated_text = apply_corrections(translated_text)
                
                # 收集翻译前后的文本
                original_texts.append(text)
                translated_texts.append(translated_text)
                if append_translation.get():
                    run.text = append_translation_to_original(text, translated_text)
                else:
                    run.text = translated_text
            # 处理图片
            for inline in run.element.xpath('.//w:drawing'):  # 解决问题的核心
                pic_elements = inline.xpath('.//a:blip/@r:embed')
                if pic_elements:  # 检查是否存在图片元素
                    pic = pic_elements[0]
                    # 保持图片在文档中的位置 and 大小
                    new_run = paragraph.add_run()
                    new_run._r.append(inline)
    
    # 处理表格
    total_tables = len(doc.tables)
    for i, table in enumerate(doc.tables, 1):
        msg = f"正在翻译 Word: 表格 {i}/{total_tables}..."
        status_label.config(text=msg)
        print(msg)
        root.update()
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        # 处理文字翻译
                        text = run.text
                        if text.strip():  # 只翻译非空文本
                            # 先获取翻译结果
                            translated_text = get_translation(text)
                            # 对翻译结果进行纠错替换
                            translated_text = apply_corrections(translated_text)
                            
                            # 收集翻译前后的文本
                            original_texts.append(text)
                            translated_texts.append(translated_text)
                            if append_translation.get():
                                run.text = append_translation_to_original(text, translated_text)
                            else:
                                run.text = translated_text
                        # 处理图片
                        for inline in run.element.xpath('.//w:drawing'):
                            pic_elements = inline.xpath('.//a:blip/@r:embed')
                            if pic_elements:  # 检查是否存在图片元素
                                pic = pic_elements[0]
                                # 保持图片在文档中的位置和大小
                                new_run = paragraph.add_run()
                                new_run._r.append(inline)
    
    # 处理插入形状中的文字
    msg = "正在处理 Word 形状文字..."
    status_label.config(text=msg)
    print(msg)
    root.update()
    for shape in doc.inline_shapes:
        try:
            # 通过 XML 访问形状中的文本
            if hasattr(shape, '_inline'):
                inline = shape._inline
                if hasattr(inline, 'graphic'):
                    graphic = inline.graphic
                    if hasattr(graphic, 'graphicData'):
                        graphic_data = graphic.graphicData
                        if hasattr(graphic_data, 'txBody'):
                            text_frame = graphic_data.txBody
                            if text_frame:
                                for paragraph in text_frame.p:
                                    for run in paragraph.r:
                                        if hasattr(run, 't'):
                                            text = run.t
                                            if text and text.strip():  # 只翻译非空文本
                                                # 先获取翻译结果
                                                translated_text = get_translation(text)
                                                # 对翻译结果进行纠错替换
                                                translated_text = apply_corrections(translated_text)
                                                
                                                # 收集翻译前后的文本
                                                original_texts.append(text)
                                                translated_texts.append(translated_text)
                                                if append_translation.get():
                                                    run.t = append_translation_to_original(text, translated_text)
                                                else:
                                                    run.t = translated_text
        except AttributeError:
            continue  # 如果无法访问相关属性，跳过处理
    
    msg = "正在保存 Word 文件..."
    status_label.config(text=msg)
    print(msg)
    root.update()
    doc.save(output_file)
    # 翻译完成后调用保存语料库函数
    save_to_corpus(original_texts, translated_texts)
    # 清空列表
    original_texts.clear()
    translated_texts.clear()
    print("保存修改后的工作表完成!")

def start_translation():
    """
    开始翻译按钮的回调函数
    """
    try:
        # 禁用按钮，显示为按下/忙碌状态
        translate_button.config(state=tk.DISABLED)
        root.update()

        input_file = input_file_entry.get()
        output_folder = output_folder_entry.get()

        if input_file and output_folder:
            msg = "开始翻译过程..."
            status_label.config(text=msg)
            print(msg)
            root.update()
            file_ext = os.path.splitext(input_file)[1]
            # 规范化输出文件夹路径
            output_folder = os.path.abspath(output_folder)

            # 检查输出文件夹是否存在，不存在则创建
            if not os.path.exists(output_folder):
                os.makedirs(output_folder)

            # 处理自定义文件名
            custom_filename = custom_filename_entry.get().strip()
            if custom_filename:
                # 保持原有的扩展名
                output_file = os.path.join(output_folder, custom_filename + file_ext)
            else:
                # 默认文件名
                output_file = os.path.join(output_folder, "translated_" + os.path.basename(input_file))

            # 在翻译开始前加载纠错语料库
            load_correction_corpus()

            # 清空之前的文本列表
            global original_texts, translated_texts
            original_texts = []
            translated_texts = []

            if input_file.lower().endswith(('.ppt', '.pptx')):
                translate_ppt(input_file, output_file)
            elif input_file.lower().endswith(('.xls', '.xlsx')):
                translate_excel(input_file, output_file)
            elif input_file.lower().endswith(('.docx')):
                translate_word(input_file, output_file)

            msg = "翻译完成！文件保存路径: " + output_file
            status_label.config(text=msg)
            #print(msg)
        else:
            msg = "请选择输入文件和输出文件夹路径。"
            status_label.config(text=msg)
            print(msg)
    except Exception as e:
        msg = f"翻译过程出错: {str(e)}"
        status_label.config(text=msg)
        print(f"错误详情: {str(e)}")
    finally:
        # 无论成功还是失败，都恢复按钮状态
        translate_button.config(state=tk.NORMAL)
        root.update()


# 定义 translation_direction 变量
translation_direction = tk.StringVar()
translation_direction.set('zh2en')  # 设置默认翻译方向为中文 -> 英文

# 输入文件选择
input_file_label = tk.Label(root, text="选择要翻译的文件:")
input_file_label.pack(pady=5)
input_file_entry = tk.Entry(root, width=50)
input_file_entry.pack(pady=5)

def select_input_file():
    file_path = filedialog.askopenfilename(
        filetypes=[
            ("Office文件", "*.docx;*.xls;*.xlsx;*.ppt;*.pptx"),
            ("Word文件", "*.docx"),
            ("Excel文件", "*.xls;*.xlsx"),
            ("PPT文件", "*.ppt;*.pptx"),
            ("所有文件", "*.*")
        ]
    )
    if file_path:
        input_file_entry.delete(0, tk.END)
        input_file_entry.insert(0, file_path)

input_file_button = tk.Button(root, text="浏览", command=select_input_file)
input_file_button.pack(pady=5)

def select_output_folder():
    folder_path = filedialog.askdirectory()
    if folder_path:
        output_folder_entry.delete(0, tk.END)
        output_folder_entry.insert(0, folder_path)

# 输出文件夹选择
output_folder_label = tk.Label(root, text="选择保存翻译后文件的文件夹:")
output_folder_label.pack(pady=5)
output_folder_entry = tk.Entry(root, width=50)
output_folder_entry.pack(pady=5)
output_folder_button = tk.Button(root, text="浏览", command=select_output_folder)
output_folder_button.pack(pady=5)

# 自定义文件名输入框
custom_filename_label = tk.Label(root, text="自定义输出文件名（可选）:")
custom_filename_label.pack(pady=5)
custom_filename_entry = tk.Entry(root, width=50)
custom_filename_entry.pack(pady=5)

# 创建一个框架用于放置翻译方向选择组件
direction_frame = tk.Frame(root)
direction_frame.pack(side=tk.RIGHT, padx=20, pady=20)

# 翻译方向选择标签，放置在框架内
translation_direction_label = tk.Label(direction_frame, text="选择翻译方向:")
translation_direction_label.pack(pady=5)

zh2ko_radio = tk.Radiobutton(direction_frame, text="中文 -> 韩文", variable=translation_direction, value='zh2ko')
zh2ko_radio.pack(pady=2)

ko2zh_radio = tk.Radiobutton(direction_frame, text="韩文 -> 中文", variable=translation_direction, value='ko2zh')
ko2zh_radio.pack(pady=2)

ko2vi_radio = tk.Radiobutton(direction_frame, text="韩文 -> 越南文", variable=translation_direction, value='ko2vi')
ko2vi_radio.pack(pady=2)

ko2en_radio = tk.Radiobutton(direction_frame, text="韩文 -> 英文", variable=translation_direction, value='ko2en')
ko2en_radio.pack(pady=2)

zh2en_radio = tk.Radiobutton(direction_frame, text="中文 -> 英文", variable=translation_direction, value='zh2en')
zh2en_radio.pack(pady=2)

en2zh_radio = tk.Radiobutton(direction_frame, text="英文 -> 中文", variable=translation_direction, value='en2zh')
en2zh_radio.pack(pady=2)

zh_tw2en_radio = tk.Radiobutton(direction_frame, text="繁体中文 -> 英文", variable=translation_direction, value='zh_tw2en')
zh_tw2en_radio.pack(pady=2)

en2zh_tw_radio = tk.Radiobutton(direction_frame, text="英文 -> 繁体中文", variable=translation_direction, value='en2zh_tw')
en2zh_tw_radio.pack(pady=2)

# 新增日文和中文翻译方向的单选按钮
zh2ja_radio = tk.Radiobutton(direction_frame, text="中文 -> 日文", variable=translation_direction, value='zh2ja')
zh2ja_radio.pack(pady=2)

ja2zh_radio = tk.Radiobutton(direction_frame, text="日文 -> 中文", variable=translation_direction, value='ja2zh')
ja2zh_radio.pack(pady=2)

# 新增英文到韩文翻译方向的单选按钮
en2ko_radio = tk.Radiobutton(direction_frame, text="英文 -> 韩文", variable=translation_direction, value='en2ko')
en2ko_radio.pack(pady=2)

# 新增中文到越南文翻译方向的单选按钮
zh2vi_radio = tk.Radiobutton(direction_frame, text="中文 -> 越南文", variable=translation_direction, value='zh2vi')
zh2vi_radio.pack(pady=2)

# 新增越南文到中文翻译方向的单选按钮
vi2zh_radio = tk.Radiobutton(direction_frame, text="越南文 -> 中文", variable=translation_direction, value='vi2zh')
vi2zh_radio.pack(pady=2)

# 添加是否在原文下方添加翻译文本的复选框
append_translation_checkbox = tk.Checkbutton(root, text="在原文下方添加翻译文本", variable=append_translation)
append_translation_checkbox.pack(pady=5)

# 添加是否生成语料库的复选框
generate_corpus_checkbox = tk.Checkbutton(root, text="生成语料库", variable=generate_corpus)
generate_corpus_checkbox.pack(pady=5)

# 语料库文件选择
corpus_file_label = tk.Label(root, text="选择纠错语料库文件:")
corpus_file_label.pack(pady=5)
corpus_file_entry = tk.Entry(root, textvariable=corpus_file_path, width=50)
corpus_file_entry.pack(pady=5)

def select_corpus_file():
    file_path = filedialog.askopenfilename(
        title="选择语料库文件",
        filetypes=[
            ("Excel文件", "*.xlsx;*.xls"),
            ("所有文件", "*.*")
        ]
    )
    if file_path:
        corpus_file_path.set(file_path)

corpus_file_button = tk.Button(root, text="浏览", command=select_corpus_file)
corpus_file_button.pack(pady=5)

# 开始翻译按钮
translate_button = tk.Button(root, text="开始翻译", command=start_translation)
translate_button.pack(pady=20)

# 状态标签
status_label = tk.Label(root, text="")
status_label.pack(pady=5)

#保存语料库
def save_to_corpus(original_texts, translated_texts):
    """
    将翻译前后的文本保存到语料库文件中
    :param original_texts: 翻译前的文本列表
    :param translated_texts: 翻译后的文本列表
    """
    if generate_corpus.get():
        msg = "正在生成语料库文件..."
        status_label.config(text=msg)
        print(msg)
        root.update()
        # 创建 Corpus 文件夹
        corpus_folder = 'Corpus'
        if not os.path.exists(corpus_folder):
            os.makedirs(corpus_folder)
        # 获取当前时间戳
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        # 根据翻译方向确定语言标签
        from_lang_label = {
            'zh2ko': 'chinese',
            'ko2zh': 'korean',
            'ko2en': 'korean',
            'zh2en': 'chinese',
            'en2zh': 'english',
            'zh_tw2en': 'chinese_traditional',
            'en2zh_tw': 'english',
            # 新增日文和中文翻译方向的语言标签
            'zh2ja': 'chinese',
            'ja2zh': 'japanese',
            'en2ko': 'english',  # 新增英文到韩文的源语言标签
            'vi2zh': 'vietnamese',  # 新增越南文到中文的源语言标签
            'zh2vi': 'chinese',
            'ko2vi': 'korean'
        }.get(translation_direction.get(), 'unknown')
        to_lang_label = {
            'zh2ko': 'korean',
            'ko2zh': 'chinese',
            'ko2en': 'english',
            'zh2en': 'english',
            'en2zh': 'chinese',
            'zh_tw2en': 'english',
            'en2zh_tw': 'chinese_traditional',
            # 新增日文和中文翻译方向的语言标签
            'zh2ja': 'japanese',
            'ja2zh': 'chinese',
            'en2ko': 'korean',  # 新增英文到韩文的目标语言标签
            'vi2zh': 'chinese',  # 新增越南文到中文的目标语言标签
            'zh2vi': 'vietnamese',
            'ko2vi': 'vietnamese'
        }.get(translation_direction.get(), 'unknown')

        # 生成 XLSX 文件名
        corpus_file = os.path.join(corpus_folder, f'Corpus_{from_lang_label}_to_{to_lang_label}_{timestamp}.xlsx')

        # 使用 pandas 创建 DataFrame 并写入 xlsx
        df = pd.DataFrame({
            '序号': range(1, len(original_texts) + 1),
            '翻译前': original_texts,
            '翻译后': translated_texts
        })
        df.to_excel(corpus_file, index=False, engine='openpyxl')

# 运行主循环
root.mainloop()
